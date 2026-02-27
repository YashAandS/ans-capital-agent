"""
modules/committee_deck.py
Committee Presentation Builder — fills the actual A&S Capital IC .pptx template
by locating shapes by name and replacing placeholder text.

The template has 4 slides:
  Slide 1: Title/Cover — property address, city/state, type, loan amount, LTV, rate
  Slide 2: Property & Loan Overview — details, financial metrics, loan breakdown table
  Slide 3: Ownership & Comps — ownership history, recorded transactions table, comps screenshot area
  Slide 4: Investment Highlights — bullet points

Uses Tavily to search Zillow for comps; Claude to generate investment highlights.
"""

import io
from datetime import datetime
from copy import deepcopy

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from anthropic import Anthropic


# ---------------------------------------------------------------------------
# Shape-name to placeholder text mapping for the IC template
# This is based on the actual shape names in AS_Capital_IC_Template (1).pptx
# ---------------------------------------------------------------------------

def _find_shape_by_name(slide, name):
    """Find a shape on a slide by its name attribute."""
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def _set_shape_text(slide, shape_name, new_text):
    """Replace the text in a named shape, preserving formatting of first run."""
    shape = _find_shape_by_name(slide, shape_name)
    if shape and shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            if para.runs:
                # Keep formatting from first run, replace text
                para.runs[0].text = str(new_text)
                # Clear other runs
                for run in para.runs[1:]:
                    run.text = ""
                return True
    return False


def _set_table_cell(slide, table_name, row, col, text):
    """Set text in a specific table cell."""
    shape = _find_shape_by_name(slide, table_name)
    if shape and shape.has_table:
        cell = shape.table.cell(row, col)
        # Preserve formatting
        if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
            cell.text_frame.paragraphs[0].runs[0].text = str(text)
        else:
            cell.text = str(text)
        return True
    return False


# ---------------------------------------------------------------------------
# Comps search via Tavily
# ---------------------------------------------------------------------------

def search_comps(tavily_api_key: str, property_address: str, property_type: str) -> list[dict]:
    """Search Zillow for comparable recent sales near the property."""
    from tavily import TavilyClient

    client = TavilyClient(api_key=tavily_api_key)
    query = (
        f"Zillow recently sold comparable properties near {property_address} "
        f"{property_type} sale price square feet bedrooms within 6 months"
    )
    results = client.search(
        query=query,
        search_depth="advanced",
        max_results=8,
        include_domains=["zillow.com"],
    )
    comps = []
    for r in results.get("results", []):
        comps.append({
            "title": r.get("title", ""),
            "url": r.get("url", ""),
            "snippet": r.get("content", ""),
        })
    return comps


def parse_comps_with_claude(api_key: str, raw_comps: list[dict], subject_address: str) -> str:
    """Use Claude to parse raw Tavily results into a clean comps summary."""
    client = Anthropic(api_key=api_key)
    comps_text = ""
    for i, c in enumerate(raw_comps, 1):
        comps_text += f"\nComp {i}:\nTitle: {c['title']}\nURL: {c['url']}\nSnippet: {c['snippet']}\n"

    response = client.messages.create(
        model="claude-4-sonnet-20250514",
        max_tokens=2048,
        system="""You are a real estate analyst at A&S Capital. Parse the search results into
a clean comparable sales summary. For each valid comp, extract: Address, Sale Price,
Sale Date, Sq Ft, Beds/Baths, Distance to subject. Then provide brief valuation commentary.
Flag comps outside the immediate neighborhood or older than 6 months.
For condos, note if at least one comp is from outside the project.""",
        messages=[{
            "role": "user",
            "content": f"Subject Property: {subject_address}\n\nRaw search results:\n{comps_text}"
        }],
    )
    return response.content[0].text


def generate_investment_highlights(api_key: str, deal_details: dict) -> str:
    """Use Claude to generate investment highlights bullet points."""
    client = Anthropic(api_key=api_key)

    detail_text = "\n".join(f"{k}: {v}" for k, v in deal_details.items() if v)

    response = client.messages.create(
        model="claude-4-opus-20250514",
        max_tokens=1024,
        system="""You are writing investment highlights for an A&S Capital investment committee
presentation. Produce 4-6 concise bullet points highlighting the key strengths of the deal.
Cover: location/market, borrower experience, loan structure, exit strategy, and any value-add.
Format each point starting with "•  " and keep each to 1-2 sentences.""",
        messages=[{
            "role": "user",
            "content": f"Generate investment highlights for this deal:\n\n{detail_text}"
        }],
    )
    return response.content[0].text


# ---------------------------------------------------------------------------
# Main deck builder
# ---------------------------------------------------------------------------

def build_committee_deck(
    template_path: str,
    anthropic_api_key: str = "",
    tavily_api_key: str = "",
    # Deal details
    property_address: str = "",
    city_state_zip: str = "",
    property_type: str = "",
    transaction_type: str = "",
    loan_type: str = "",
    loan_number: str = "",
    # Financial metrics
    total_loan_amount: float = 0,
    ltv_to_arv: float = 0,
    interest_rate: float = 0,
    purchase_price: float = 0,
    purchase_date: str = "",
    as_is_value: float = 0,
    after_repair_value: float = 0,
    rehab_budget: float = 0,
    # Loan breakdown
    initial_loan: float = 0,
    initial_ltc: float = 0,
    interest_reserve: float = 0,
    holdback_rehab: float = 0,
    holdback_ltc: float = 0,
    total_ltc: float = 0,
    # Property details
    year_built: str = "",
    square_footage: str = "",
    num_buildings: str = "",
    lot_sf: str = "",
    lot_acres: str = "",
    subdivision: str = "",
    loan_term: str = "",
    # Ownership
    assessment_owner: str = "",
    buyer_entity: str = "",
    purchase_history: str = "",
    existing_financing: str = "",
    # Recorded transactions (list of dicts)
    recorded_transactions: list = None,
    # Investment highlights
    investment_highlights: str = "",
    # Classification
    classification: str = "",
) -> io.BytesIO:
    """
    Fill the A&S Capital IC template and return as BytesIO.
    """
    prs = Presentation(template_path)
    today = datetime.now()

    # ===== SLIDE 1: Title/Cover =====
    slide1 = prs.slides[0]

    _set_shape_text(slide1, "Text 4", property_address or "[Property Address]")
    _set_shape_text(slide1, "Text 5", city_state_zip or "[City, State ZIP]")

    # Property type / transaction / loan type line
    info_line = "   ·   ".join(filter(None, [property_type, transaction_type, loan_type]))
    _set_shape_text(slide1, "Text 7", info_line or "[Property Type]   ·   [Transaction Type]   ·   [Loan Type]")

    # Footer with loan # and date
    footer = f"A&S Capital LLC   ·   Loan #{loan_number or '____'}   ·   Prepared {today.strftime('%B %Y')}"
    _set_shape_text(slide1, "Text 9", footer)

    # Loan amount
    if total_loan_amount:
        _set_shape_text(slide1, "Text 10", f"${total_loan_amount:,.0f}")

    # LTV to ARV
    if ltv_to_arv:
        _set_shape_text(slide1, "Text 12", f"{ltv_to_arv:.0%}" if ltv_to_arv < 1 else f"{ltv_to_arv:.1f}%")
    _set_shape_text(slide1, "Text 13", "LTV to ARV")

    # Interest rate
    if interest_rate:
        _set_shape_text(slide1, "Text 14", f"{interest_rate:.2%}" if interest_rate < 1 else f"{interest_rate:.1f}%")
    _set_shape_text(slide1, "Text 15", "Interest Rate")

    # ===== SLIDE 2: Property & Loan Overview =====
    slide2 = prs.slides[1]

    # Address block
    _set_shape_text(slide2, "Text 6", f"{property_address}, {city_state_zip}" if property_address else "[Street Address, City, State ZIP]")

    # Type & classification
    type_class = "   ·   ".join(filter(None, [property_type, classification]))
    _set_shape_text(slide2, "Text 9", type_class or "[Property Type]  ·  [Classification]")

    # Year built / area
    area_parts = [year_built, f"{square_footage} SF" if square_footage else None, f"{num_buildings} Bldg" if num_buildings else None]
    _set_shape_text(slide2, "Text 12", "  ·  ".join(filter(None, area_parts)) or "[Year]  ·  [SF]  ·  [# Buildings]")

    # Lot
    lot_parts = [f"{lot_sf} SF" if lot_sf else None, f"({lot_acres} AC)" if lot_acres else None, subdivision]
    _set_shape_text(slide2, "Text 15", "  ·  ".join(filter(None, lot_parts)) or "[Lot SF] ([AC])  ·  [Subdivision]")

    # Transaction line
    txn_line = "  ·  ".join(filter(None, [transaction_type, f"{loan_term} Bridge Loan" if loan_term else None]))
    _set_shape_text(slide2, "Text 18", txn_line or "[Transaction Type]  ·  [Loan Term] Bridge Loan")

    # Financial metrics boxes
    if purchase_price:
        _set_shape_text(slide2, "Text 22", f"${purchase_price:,.0f}")
    _set_shape_text(slide2, "Text 23", "Purchase Price")
    _set_shape_text(slide2, "Text 24", purchase_date or "[Date]")

    if as_is_value:
        _set_shape_text(slide2, "Text 27", f"${as_is_value:,.0f}")
    _set_shape_text(slide2, "Text 28", "As-Is Value")

    if after_repair_value:
        _set_shape_text(slide2, "Text 32", f"${after_repair_value:,.0f}")
    _set_shape_text(slide2, "Text 33", "After Repair Value")

    if rehab_budget:
        _set_shape_text(slide2, "Text 37", f"${rehab_budget:,.0f}")
    _set_shape_text(slide2, "Text 38", "Rehab Budget")

    # Loan breakdown table (Table 0 on slide 2)
    if initial_loan:
        _set_table_cell(slide2, "Table 0", 1, 1, f"${initial_loan:,.0f}")
        _set_table_cell(slide2, "Table 0", 1, 2, f"{initial_ltc:.0%}" if initial_ltc else "—")
    if interest_reserve:
        _set_table_cell(slide2, "Table 0", 2, 1, f"${interest_reserve:,.0f}")
    if holdback_rehab:
        _set_table_cell(slide2, "Table 0", 3, 1, f"${holdback_rehab:,.0f}")
        _set_table_cell(slide2, "Table 0", 3, 2, f"{holdback_ltc:.0%}" if holdback_ltc else "—")
    if total_loan_amount:
        _set_table_cell(slide2, "Table 0", 4, 1, f"${total_loan_amount:,.0f}")
        _set_table_cell(slide2, "Table 0", 4, 2, f"{total_ltc:.0%}" if total_ltc else "—")

    # ===== SLIDE 3: Ownership & Comps =====
    slide3 = prs.slides[2]

    _set_shape_text(slide3, "Text 7", assessment_owner or "Current Owner Name")
    _set_shape_text(slide3, "Text 9", buyer_entity or "Current Entity Name")

    if purchase_price and purchase_date:
        _set_shape_text(slide3, "Text 11", f"{purchase_date}  ·  Price: ${purchase_price:,.0f}")
    _set_shape_text(slide3, "Text 13", existing_financing or "$[___] via [Lender Name]")

    # Recorded transactions table (Table 0 on slide 3)
    if recorded_transactions:
        for i, txn in enumerate(recorded_transactions[:5]):
            row = i + 1
            _set_table_cell(slide3, "Table 0", row, 0, txn.get("address", f"{i+1}."))
            _set_table_cell(slide3, "Table 0", row, 1, txn.get("date", ""))
            _set_table_cell(slide3, "Table 0", row, 2, txn.get("price", ""))

    # Comps area — add comps text to the instruction box
    if tavily_api_key and anthropic_api_key and property_address:
        try:
            raw_comps = search_comps(tavily_api_key, property_address, property_type)
            if raw_comps:
                comps_text = parse_comps_with_claude(anthropic_api_key, raw_comps, property_address)
                _set_shape_text(slide3, "Text 21", comps_text)
        except Exception as e:
            _set_shape_text(slide3, "Text 21", f"Comps search error: {e}\nAdd comps manually.")

    # ===== SLIDE 4: Investment Highlights =====
    slide4 = prs.slides[3]

    highlights = investment_highlights
    if not highlights and anthropic_api_key:
        try:
            details = {
                "property_address": property_address,
                "city_state_zip": city_state_zip,
                "property_type": property_type,
                "loan_type": loan_type,
                "loan_amount": f"${total_loan_amount:,.0f}" if total_loan_amount else "",
                "as_is_value": f"${as_is_value:,.0f}" if as_is_value else "",
                "arv": f"${after_repair_value:,.0f}" if after_repair_value else "",
                "rehab_budget": f"${rehab_budget:,.0f}" if rehab_budget else "",
                "interest_rate": f"{interest_rate:.2%}" if interest_rate else "",
                "borrower": buyer_entity,
                "transaction_type": transaction_type,
            }
            highlights = generate_investment_highlights(anthropic_api_key, details)
        except Exception:
            highlights = "Enter investment highlights here...\n\n•  \n\n•  \n\n•  \n\n•  "

    _set_shape_text(slide4, "Text 5", highlights)

    # Write to BytesIO
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output
