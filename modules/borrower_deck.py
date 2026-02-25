"""
modules/borrower_deck.py
Borrower Presentation Builder — creates a clean borrower-facing loan proposal
using python-pptx (generates from scratch, not from IC template, since the
borrower deck has a different structure than the committee deck).

Slides:
  1. Title/Cover
  2. Loan Terms Summary
  3. Property Details
  4. Estimated Costs & Fees
  5. Timeline & Process
  6. Requirements & Next Steps
  7. Contact Information
"""

import io
from datetime import datetime, timedelta

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ---------------------------------------------------------------------------
# Styling constants (matching A&S Capital branding)
# ---------------------------------------------------------------------------
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE = RGBColor(0xD6, 0xE4, 0xF0)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
ACCENT_GOLD = RGBColor(0xC6, 0x8A, 0x00)


def _add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY

    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = LIGHT_BLUE
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(12)

    # Tagline
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.4))
    tf2 = txBox2.text_frame
    p3 = tf2.paragraphs[0]
    p3.text = "Private Lending Solutions  |  Bridge  |  Fix & Flip  |  DSCR  |  Construction"
    p3.font.size = Pt(9)
    p3.font.color.rgb = LIGHT_BLUE
    p3.alignment = PP_ALIGN.CENTER
    return slide


def _add_content_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    tf.margin_left = Inches(0.4)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return slide


def _add_kv_block(slide, items, left=0.4, top=0.9, width=9.0, font_size=10):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(4))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (key, value) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run_k = p.add_run()
        run_k.text = f"{key}:  "
        run_k.font.size = Pt(font_size)
        run_k.font.bold = True
        run_k.font.color.rgb = NAVY
        run_v = p.add_run()
        run_v.text = str(value)
        run_v.font.size = Pt(font_size)
        run_v.font.color.rgb = DARK_GRAY
        p.space_after = Pt(5)


def _add_numbered_list(slide, items, left=0.4, top=0.9, width=9.0, font_size=10):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(4))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"  {i + 1}.  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(4)


def build_borrower_deck(
    # Deal details
    property_address: str = "",
    city_state_zip: str = "",
    property_type: str = "",
    num_units: int = 1,
    square_footage: int = 0,
    as_is_value: float = 0,
    arv: float = 0,
    borrower_name: str = "",
    entity_name: str = "",
    loan_purpose: str = "Purchase",
    loan_type: str = "RTL",
    loan_amount: float = 0,
    loan_term_months: int = 12,
    interest_rate: float = 0.10,
    origination_fee_pct: float = 0.02,
    rehab_budget: float = 0,
    interest_reserve_months: int = 0,
    additional_requirements: str = "",
    contact_name: str = "A&S Capital Originations",
    contact_email: str = "originations@ascapital.com",
    contact_phone: str = "305.749.0848",
) -> io.BytesIO:
    """Build a borrower-facing loan proposal presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Emu(5143500)  # Match IC template height

    today = datetime.now()
    today_str = today.strftime("%B %d, %Y")
    effective_arv = arv if arv else as_is_value
    ltv = loan_amount / as_is_value if as_is_value > 0 else 0
    origination_fee = loan_amount * origination_fee_pct
    monthly_interest = loan_amount * interest_rate / 12
    interest_reserve = monthly_interest * interest_reserve_months
    net_proceeds = loan_amount - origination_fee - interest_reserve

    # Slide 1: Title
    _add_title_slide(prs, "A&S CAPITAL",
        f"Loan Proposal for {borrower_name}\n{property_address}\n{today_str}")

    # Slide 2: Loan Terms
    slide = _add_content_slide(prs, "Loan Terms Summary")
    _add_kv_block(slide, [
        ("Loan Amount", f"${loan_amount:,.0f}"),
        ("Loan Purpose", loan_purpose),
        ("Loan Type", loan_type),
        ("Term", f"{loan_term_months} months"),
        ("Interest Rate", f"{interest_rate:.2%} (interest only)"),
        ("Origination Fee", f"{origination_fee_pct:.1%} (${origination_fee:,.0f})"),
        ("LTV", f"{ltv:.1%}"),
        ("Rehab Budget", f"${rehab_budget:,.0f}" if rehab_budget > 0 else "N/A"),
        ("Monthly Payment (est.)", f"${monthly_interest:,.0f}"),
        ("Prepayment Penalty", "None"),
    ])

    # Slide 3: Property Details
    slide = _add_content_slide(prs, "Property Details")
    _add_kv_block(slide, [
        ("Address", property_address),
        ("City / State / ZIP", city_state_zip),
        ("Property Type", property_type),
        ("Units", str(num_units)),
        ("Square Footage", f"{square_footage:,}" if square_footage else "TBD"),
        ("As-Is Value", f"${as_is_value:,.0f}"),
        ("After-Repair Value", f"${effective_arv:,.0f}" if arv else "N/A"),
    ])

    # Slide 4: Costs & Fees
    slide = _add_content_slide(prs, "Estimated Costs & Fees")
    _add_kv_block(slide, [
        ("Origination Fee", f"${origination_fee:,.0f}  ({origination_fee_pct:.1%})"),
        ("Interest Reserve", f"${interest_reserve:,.0f}  ({interest_reserve_months} months)"),
        ("Est. Total Closing Costs", f"${origination_fee + interest_reserve:,.0f}"),
        ("Net Loan Proceeds", f"${net_proceeds:,.0f}"),
    ], top=0.9)

    # Disclaimer
    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(3.0), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ("Note: Closing costs are estimates. Actual costs may include appraisal, title insurance, "
              "recording fees, legal fees, and other third-party costs.")
    p.font.size = Pt(8)
    p.font.italic = True
    p.font.color.rgb = MEDIUM_GRAY

    # Slide 5: Timeline
    slide = _add_content_slide(prs, "Timeline & Process")
    _add_numbered_list(slide, [
        f"Application & Term Sheet — {today_str}",
        f"Signed Term Sheet & Deposit — by {(today + timedelta(days=3)).strftime('%B %d, %Y')}",
        f"Appraisal Ordered — by {(today + timedelta(days=5)).strftime('%B %d, %Y')}",
        f"Underwriting & Title Review — {(today + timedelta(days=5)).strftime('%b %d')} to {(today + timedelta(days=14)).strftime('%b %d')}",
        f"Loan Approval / Committee — by {(today + timedelta(days=15)).strftime('%B %d, %Y')}",
        f"Clear to Close — by {(today + timedelta(days=18)).strftime('%B %d, %Y')}",
        f"Targeted Closing — {(today + timedelta(days=21)).strftime('%B %d, %Y')}",
    ])

    # Slide 6: Requirements
    slide = _add_content_slide(prs, "Requirements & Next Steps")
    reqs = [
        "Signed Term Sheet and application",
        "Good faith deposit (refundable at closing)",
        "Entity documents (Operating Agreement, Certificate of Good Standing)",
        "Personal financial statement and bank statements (2 months)",
        "Proof of insurance (or binder prior to closing)",
        "Schedule of Real Estate Owned (REO)",
        "Photo ID for all guarantors",
        "Purchase contract or payoff statement (if refinance)",
    ]
    if rehab_budget > 0:
        reqs.append("Detailed scope of work with contractor bids")
    if additional_requirements:
        for r in additional_requirements.strip().split("\n"):
            r = r.strip()
            if r:
                reqs.append(r)
    _add_numbered_list(slide, reqs)

    # Slide 7: Contact
    slide = _add_content_slide(prs, "Contact Information")
    txBox = slide.shapes.add_textbox(Inches(2), Inches(1.5), Inches(6), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "A&S CAPITAL LLC"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER

    for label, val in [("Contact", contact_name), ("Email", contact_email), ("Phone", contact_phone)]:
        if val:
            p2 = tf.add_paragraph()
            run_k = p2.add_run()
            run_k.text = f"{label}: "
            run_k.font.size = Pt(11)
            run_k.font.bold = True
            run_k.font.color.rgb = DARK_GRAY
            run_v = p2.add_run()
            run_v.text = val
            run_v.font.size = Pt(11)
            run_v.font.color.rgb = DARK_GRAY
            p2.alignment = PP_ALIGN.CENTER
            p2.space_after = Pt(4)

    # Disclaimer
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(0.8))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = ("This Financing Quote is for discussion purposes only and does not constitute a loan "
              "approval or commitment. Terms subject to change after complete application review.")
    p.font.size = Pt(7)
    p.font.italic = True
    p.font.color.rgb = MEDIUM_GRAY
    p.alignment = PP_ALIGN.CENTER

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output
